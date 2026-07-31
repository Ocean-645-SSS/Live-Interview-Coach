"""文档解析器：从 PDF、DOCX、PPTX、XLSX 和纯文本文件中提取文本。
这部分逻辑参考 HKUDS/LightRAG 的 document_routes.py 流程，
并改造成当前 RAG Core Service 使用的同步解析函数。
"""

from __future__ import annotations

from io import BytesIO
from zipfile import BadZipFile

from docx import Document
from docx.opc.exceptions import OpcError
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph
from lxml.etree import XMLSyntaxError
from openpyxl import load_workbook
from openpyxl.utils.exceptions import InvalidFileException
from pptx import Presentation
from pptx.exc import PythonPptxError
from pypdf import PdfReader
from pypdf.errors import PyPdfError

# ---------------------------------------------------------------------------
# UTF-8 文本文件扩展名。
# ---------------------------------------------------------------------------
TEXT_EXTENSIONS: set[str] = {
    ".txt",
    ".md",
    ".mdx",
    ".html",
    ".htm",
    ".tex",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".csv",
    ".log",
    ".conf",
    ".ini",
    ".properties",
    ".sql",
    ".bat",
    ".sh",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".py",
    ".java",
    ".js",
    ".ts",
    ".swift",
    ".go",
    ".rb",
    ".php",
    ".css",
    ".scss",
    ".less",
}


# ---------------------------------------------------------------------------
# 分发器
# ---------------------------------------------------------------------------

# 二进制文件扩展名会路由到对应解析器，PDF 固定在这里处理。
BINARY_EXTENSIONS: dict[str,str] ={
    ".pdf":"pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
}

#合并TEXT_EXTENSIONS+BINARY_EXTENSIONS得到完整的合法文件后缀
ALL_SUPPORTED_EXTENSIONS=TEXT_EXTENSIONS | set(BINARY_EXTENSIONS)

def parse_file_content(file_bytes:bytes,extension:str,**kwargs:object)->str:
    """将文件解析为string类型，分别处理text类型和binary类型"""
    if not isinstance(file_bytes, bytes):
        raise TypeError("file_bytes 必须是 bytes")
    if not isinstance(extension, str):
        raise TypeError("extension 必须是 str")

    ext = extension.strip().lower()
    if not ext.startswith("."):
        ext = f".{ext}"

    #先处理文本文件（.txt/.md....）
    if ext in TEXT_EXTENSIONS:
        try:
            text=file_bytes.decode("utf-8")
        except UnicodeDecodeError as exc:
            raise ValueError(
                "当前文件不是合法 UTF-8，请转换为 UTF-8 格式后再处理"
            ) from exc
        if text.strip() == "":
            raise ValueError("文本内容为空！")
        if text.startswith("b'") or text.startswith('b"'):
            raise ValueError("文件看起来包含二进制数据")
        return text

    #处理二进制类型文件(.pdf/.docs....)
    dispatch=BINARY_EXTENSIONS.get(ext)
    if dispatch is None:
        raise ValueError(
            f"不支持当前扩展名:{ext}"
            f"支持的扩展名有：{sorted(ALL_SUPPORTED_EXTENSIONS)}"
        )
    if dispatch=='pdf':
        password=kwargs.get("password")
        pdf_password = str(password) if password is not None else None
        return extract_pdf(file_bytes, password=pdf_password)
    if dispatch == "docx":
        return extract_docx(file_bytes)
    if dispatch == "pptx":
        return extract_pptx(file_bytes)
    if dispatch == "xlsx":
        return extract_xlsx(file_bytes)

    raise ValueError(f"未处理的二进制扩展名: {ext}")  # pragma: no cover

# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------
def extract_pdf(file_bytes: bytes, password: str | None = None) -> str:
    """优先使用 Docling 解析 PDF，不可用或解析失败时回退到 pypdf。

    Docling 不支持通过这里的内存接口传入 PDF 密码，因此带密码的
    PDF 直接交给 pypdf 处理。
    """

    if password is None:
        try:
            return _extract_pdf_with_docling(file_bytes)
        except Exception:
            # Docling 是可选增强能力。导入、模型初始化或转换失败时，
            # 使用稳定且轻量的 pypdf 路径继续解析。
            pass

    return _extract_pdf_with_pypdf(file_bytes, password=password)


def _extract_pdf_with_docling(file_bytes: bytes) -> str:
    """使用延迟导入的 Docling 从内存 PDF 中提取 Markdown 文本。"""

    # 延迟导入避免未安装 Docling 时影响整个文档解析模块启动。
    from docling.datamodel.base_models import DocumentStream
    from docling.document_converter import DocumentConverter

    source = DocumentStream(name="uploaded.pdf", stream=BytesIO(file_bytes))
    result = DocumentConverter().convert(source)
    text = result.document.export_to_markdown().strip()

    if not text:
        raise ValueError("Docling 未从 PDF 中提取出文字")
    return text


def _extract_pdf_with_pypdf(
    file_bytes: bytes,
    password: str | None = None,
) -> str:
    """使用 pypdf 解析 PDF，并处理加密文件。"""

    #转换成内存中PDF文件->转换为可读取的PDF文件
    try:
        reader=PdfReader(BytesIO(file_bytes))
    except (PyPdfError, OSError, ValueError) as exc:
        raise ValueError("PDF 文件可能被损坏或格式无效") from exc

    #需要解密
    if reader.is_encrypted:
        try:
            decrypt_result=reader.decrypt(password or "")
        except (PyPdfError, OSError, ValueError) as exc:
            raise ValueError("PDF 解密失败") from exc

        if decrypt_result==0:
            if password is not None:
                raise ValueError("PDF 密码错误")
            raise ValueError("PDF 已加密，但没有提供密码")

    #解析出文字
    parts:list[str]=[]
    try:
        for page in reader.pages:
            text=page.extract_text()
            if text and text.strip():
                parts.append(text.strip())
    except (PyPdfError, OSError, ValueError) as exc:
        raise ValueError("PDF 文字提取失败") from exc

    text = "\n".join(parts).strip()
    if not text:
        raise ValueError("PDF 中没有可提取的文字")

    return text

# ---------------------------------------------------------------------------
# DOCS
# ---------------------------------------------------------------------------
_OLE_COMPOUND_FILE_SIGNATURE = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
_OFFICE_PARSE_ERRORS = (
    BadZipFile,
    InvalidFileException,
    KeyError,
    OSError,
    OpcError,
    PythonPptxError,
    ValueError,
    XMLSyntaxError,
)


def _reject_encrypted_office_file(file_bytes: bytes, extension: str) -> None:
    """拒绝采用 OLE 加密容器保存的 OOXML 文档。"""

    if file_bytes.startswith(_OLE_COMPOUND_FILE_SIGNATURE):
        raise ValueError(f"{extension.upper()} 文件已加密，暂不支持解析")


def extract_docx(file_bytes:bytes)->str:
    """解析 DOCX 文件中的文本和表格"""

    _reject_encrypted_office_file(file_bytes, "DOCX")
    try:
        docx_file=Document(BytesIO(file_bytes))
    except _OFFICE_PARSE_ERRORS as exc:
        raise ValueError("DOCX 文件可能被损坏或格式无效") from exc
    content_parts:list[str]=[]

    try:
        for element in docx_file.element.body:
            # 元素是普通段落
            if element.tag.endswith("p"):
                paragraph = Paragraph(element, docx_file)
                if paragraph.text.strip():
                    content_parts.append(paragraph.text.strip())

            # 当前元素是表格
            elif element.tag.endswith("tbl"):
                table = DocxTable(element, docx_file)

                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        content_parts.append("\t".join(cells))
    except _OFFICE_PARSE_ERRORS as exc:
        raise ValueError("DOCX 内容提取失败") from exc

    text = "\n".join(content_parts).strip()
    if not text:
        raise ValueError("DOCX 中没有可提取的文字")
    return text

# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------
def extract_pptx(file_bytes:bytes)->str:
    """解析 PPTX 图片中的文字"""

    _reject_encrypted_office_file(file_bytes, "PPTX")
    try:
        reader=Presentation(BytesIO(file_bytes))
    except _OFFICE_PARSE_ERRORS as exc:
        raise ValueError("PPTX 文件可能被损坏或格式无效") from exc

    parts:list[str]=[]
    try:
        for slide in reader.slides:
            slide_parts:list[str]=[]
            for shape in slide.shapes:
                slide_parts.extend(_extract_shape_text(shape))
            if slide_parts:
                parts.extend(slide_parts)

    except _OFFICE_PARSE_ERRORS as exc:
        raise ValueError("PPTX 文字提取失败") from exc

    text = "\n".join(parts).strip()
    if not text:
        raise ValueError("PPTX 中没有可提取的文字")

    return text

def _extract_shape_text(shape) -> list[str]:
    """提取一个 PPTX 图形中的文字。"""

    from pptx.enum.shapes import MSO_SHAPE_TYPE

    parts: list[str] = []

    # 1. 普通文本框、标题、占位符
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            parts.append(text)

    # 2. 表格
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]

            if any(cells):
                parts.append("\t".join(cells))

    # 3. 分组图形：递归处理组内每个子图形
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            parts.extend(_extract_shape_text(child_shape))

    # 4. 图表标题
    if getattr(shape, "has_chart", False):
        chart = shape.chart

        if chart.has_title:
            title = chart.chart_title.text_frame.text.strip()
            if title:
                parts.append(title)

    return parts
# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------
def extract_xlsx(file_bytes:bytes)->str:
    """读取 Excel 工作簿，把每个工作表转换成纯文本，单元格之间用\t隔开"""

    _reject_encrypted_office_file(file_bytes, "XLSX")
    try:
        workbook=load_workbook(BytesIO(file_bytes),read_only=True,data_only=True)
    except _OFFICE_PARSE_ERRORS as exc:
        raise ValueError("XLSX 文件可能被损坏或格式无效") from exc

    content_parts:list[str]=[]
    has_content = False

    try:
        #对于每一张工作表
        for sheet in workbook.worksheets:
            sheet_rows:list[str]=[]

            #读取当前工作表中每一行
            for row in sheet.iter_rows(values_only=True):
                values = [
                    "" if value is None else str(value).strip()
                    for value in row
                ]

                #除去行末没有意义的空单元格
                while values and values[-1]=="":
                    values.pop()

                if values:
                    has_content=True
                    sheet_rows.append("\t".join(values))

            if sheet_rows:
                content_parts.append(f"Sheet: {sheet.title}")
                content_parts.extend(sheet_rows)

    except _OFFICE_PARSE_ERRORS as exc:
        raise ValueError("XLSX 内容提取失败") from exc
    finally:
        workbook.close()

    if not has_content:
        raise ValueError("XLSX 中没有可提取的内容")

    return "\n".join(content_parts)


