"""测试 M1 最小文档解析器的输入边界。"""

from io import BytesIO

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

from liverag.rag import doc_parser
from liverag.rag.doc_parser import parse_file_content


def _pdf_with_text(text: str = "LiveRAG PDF") -> bytes:
    """生成含基础 Type1 字体文本的最小 PDF。"""

    buffer = BytesIO()
    writer = PdfWriter()
    page = writer.add_blank_page(width=200, height=100)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    page[NameObject("/Resources")] = DictionaryObject(
        {
            NameObject("/Font"): DictionaryObject(
                {NameObject("/F1"): writer._add_object(font)}
            )
        }
    )
    stream = DecodedStreamObject()
    stream.set_data(f"BT /F1 12 Tf 20 50 Td ({text}) Tj ET".encode("ascii"))
    page[NameObject("/Contents")] = writer._add_object(stream)
    writer.write(buffer)
    return buffer.getvalue()


def test_parse_utf8_txt_returns_decoded_content():
    """测试正常UTF-8和.txt文件可以被解析，并且和原文一致"""
    content = "LiveRAG 支持 UTF-8 文本。\n第二行内容。"

    result = parse_file_content(content.encode("utf-8"), ".txt")

    assert result == content


def test_parse_markdown_returns_source_unchanged():
    """测试正常.md文件可以被解析，并且和原文一致"""
    content = "# LiveRAG\n\n- 保留 Markdown 标记\n- 保留原始换行\n"

    result = parse_file_content(content.encode("utf-8"), ".md")

    assert result == content


def test_extension_is_normalized():
    assert parse_file_content(b"content", " TXT ") == "content"


def test_parse_empty_file_raises_value_error():
    """空文件会被拒绝"""
    with pytest.raises(ValueError):
        parse_file_content(b"", ".txt")


def test_parse_invalid_utf8_raises_value_error():
    """验证无法按UTF-8解码的字节会抛出ValueError"""
    with pytest.raises(ValueError):
        parse_file_content(b"\xff\xfe\xfa", ".txt")


def test_parse_unsupported_extension_raises_value_error():
    """损坏的 PDF 会转换成稳定的 ValueError。"""
    with pytest.raises(ValueError):
        parse_file_content(b"not a real PDF", ".pdf")


def test_parse_pdf_extracts_text(monkeypatch):
    """Docling 不可用时，pypdf 回退路径仍能提取真实 PDF 文本。"""

    def unavailable_docling(_):
        raise ModuleNotFoundError("docling")

    monkeypatch.setattr(doc_parser, "_extract_pdf_with_docling", unavailable_docling)

    assert "LiveRAG PDF" in parse_file_content(_pdf_with_text(), ".pdf")


def test_parse_unsupported_extension_rejected():
    with pytest.raises(ValueError, match="不支持当前扩展名"):
        parse_file_content(b"content", ".exe")


def test_parse_docx_extracts_paragraph_and_table():
    buffer = BytesIO()
    document = Document()
    document.add_paragraph("段落内容")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "姓名"
    table.cell(0, 1).text = "小明"
    document.save(buffer)

    result = parse_file_content(buffer.getvalue(), ".docx")

    assert "段落内容" in result
    assert "姓名\t小明" in result
    assert result.index("段落内容") < result.index("姓名\t小明")


def test_parse_empty_docx_raises_value_error():
    buffer = BytesIO()
    Document().save(buffer)

    with pytest.raises(ValueError, match="没有可提取"):
        parse_file_content(buffer.getvalue(), ".docx")


@pytest.mark.parametrize("extension", [".docx", ".pptx", ".xlsx"])
def test_encrypted_office_file_has_stable_error(extension):
    encrypted_container = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"encrypted"

    with pytest.raises(ValueError, match="文件已加密，暂不支持解析"):
        parse_file_content(encrypted_container, extension)


@pytest.mark.parametrize("extension", [".docx", ".pptx", ".xlsx"])
def test_corrupt_office_file_has_stable_error(extension):
    with pytest.raises(ValueError, match=r"损坏|格式无效"):
        parse_file_content(b"not a valid office package", extension)


def test_parse_pptx_extracts_text():
    buffer = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "演示标题"
    presentation.save(buffer)

    assert "演示标题" in parse_file_content(buffer.getvalue(), ".pptx")


def test_parse_pptx_extracts_table():
    buffer = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, 0, 0, 1_000_000, 1_000_000).table
    table.cell(0, 0).text = "姓名"
    table.cell(0, 1).text = "年龄"
    table.cell(1, 0).text = "小明"
    table.cell(1, 1).text = "20"
    presentation.save(buffer)

    result = parse_file_content(buffer.getvalue(), ".pptx")

    assert "姓名\t年龄" in result
    assert "小明\t20" in result


def test_parse_empty_pptx_raises_value_error():
    buffer = BytesIO()
    Presentation().save(buffer)

    with pytest.raises(ValueError, match="没有可提取"):
        parse_file_content(buffer.getvalue(), ".pptx")


def test_parse_xlsx_extracts_cells():
    buffer = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "学生"
    sheet.append(["姓名", "年龄"])
    sheet.append(["小明", 20])
    workbook.save(buffer)

    result = parse_file_content(buffer.getvalue(), ".xlsx")

    assert "Sheet: 学生" in result
    assert "姓名\t年龄" in result
    assert "小明\t20" in result


def test_parse_xlsx_preserves_multiple_sheets_and_middle_empty_cell():
    buffer = BytesIO()
    workbook = Workbook()
    first = workbook.active
    first.title = "第一页"
    first.append(["姓名", None, "城市"])
    second = workbook.create_sheet("第二页")
    second.append(["小明", 20])
    workbook.save(buffer)

    result = parse_file_content(buffer.getvalue(), ".xlsx")

    assert "Sheet: 第一页" in result
    assert "姓名\t\t城市" in result
    assert "Sheet: 第二页" in result
    assert "小明\t20" in result


def test_parse_empty_xlsx_raises_value_error():
    buffer = BytesIO()
    Workbook().save(buffer)

    with pytest.raises(ValueError, match="没有可提取"):
        parse_file_content(buffer.getvalue(), ".xlsx")


def test_parse_empty_pdf_raises_value_error():
    buffer = BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.write(buffer)

    with pytest.raises(ValueError, match="没有可提取"):
        parse_file_content(buffer.getvalue(), ".pdf")


def test_encrypted_pdf_distinguishes_missing_and_wrong_password():
    buffer = BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.encrypt("correct-password")
    writer.write(buffer)
    encrypted_pdf = buffer.getvalue()

    with pytest.raises(ValueError, match="没有提供密码"):
        parse_file_content(encrypted_pdf, ".pdf")

    with pytest.raises(ValueError, match="密码错误"):
        parse_file_content(encrypted_pdf, ".pdf", password="wrong-password")

    reader = PdfReader(BytesIO(encrypted_pdf), password="correct-password")
    assert reader.is_encrypted


def test_pdf_prefers_docling(monkeypatch):
    monkeypatch.setattr(
        doc_parser,
        "_extract_pdf_with_docling",
        lambda _: "Docling result",
    )

    def unexpected_pypdf(*args, **kwargs):
        raise AssertionError("Docling 成功时不应调用 pypdf")

    monkeypatch.setattr(doc_parser, "_extract_pdf_with_pypdf", unexpected_pypdf)

    assert doc_parser.extract_pdf(b"pdf") == "Docling result"


def test_pdf_falls_back_to_pypdf_when_docling_fails(monkeypatch):
    def failed_docling(_):
        raise ModuleNotFoundError("docling")

    monkeypatch.setattr(doc_parser, "_extract_pdf_with_docling", failed_docling)
    monkeypatch.setattr(
        doc_parser,
        "_extract_pdf_with_pypdf",
        lambda file_bytes, password=None: "pypdf result",
    )

    assert doc_parser.extract_pdf(b"pdf") == "pypdf result"
