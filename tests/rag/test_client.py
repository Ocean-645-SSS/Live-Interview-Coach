"""验证 RAG Core HTTP API 的统一响应 envelope。"""

from io import BytesIO
from pathlib import Path
from typing import Any
from uuid import UUID

import pytest
from docx import Document
from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

import liverag.rag.server as server
from liverag.rag import doc_parser
from liverag.rag.knowledge_base import KnowledgeBaseStore
from liverag.rag.metadata_store import MetadataStore


def _pdf_bytes(text: str, *, password: str | None = None) -> bytes:
    buffer = BytesIO()
    writer = PdfWriter()
    page = writer.add_blank_page(width=200, height=100)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    page[NameObject("/Resources")] = DictionaryObject(
        {
            NameObject("/Font"): DictionaryObject(
                {NameObject("/F1"): writer._add_object(font)}
            )
        }
    )
    stream = DecodedStreamObject()
    stream.set_data(f"BT /F1 12 Tf 20 50 Td ({text}) Tj ET".encode("ascii"))
    page[NameObject("/Contents")] = writer._add_object(stream)
    if password is not None:
        writer.encrypt(password)
    writer.write(buffer)
    return buffer.getvalue()


def _docx_bytes(text: str) -> bytes:
    buffer = BytesIO()
    document = Document()
    document.add_paragraph(text)
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(text: str) -> bytes:
    buffer = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = text
    presentation.save(buffer)
    return buffer.getvalue()


def _xlsx_bytes(text: str) -> bytes:
    buffer = BytesIO()
    workbook = Workbook()
    workbook.active.append([text])
    workbook.save(buffer)
    return buffer.getvalue()


class FakeHttpRagEngine:
    """保存当前 KB 文本并模拟 LightRAG 的索引、job 和查询结果。"""

    def __init__(self, kb_id: str) -> None:
        self.kb_id = kb_id
        self.documents: dict[str, dict[str, str]] = {}
        self.jobs: dict[str, list[str]] = {}

    async def enqueue_documents(
        self,
        *,
        texts: list[str],
        file_sources: list[str],
        document_ids: list[str],
        track_id: str,
    ) -> None:
        for text, source, document_id in zip(
            texts,
            file_sources,
            document_ids,
            strict=True,
        ):
            self.documents[document_id] = {"text": text, "source": source}
        self.jobs[track_id] = list(document_ids)

    async def job(self, job_id: str) -> dict[str, Any]:
        document_ids = self.jobs.get(job_id)
        if document_ids is None:
            raise KeyError(f"job not found: {job_id}")
        return {
            "job_id": job_id,
            "kb_id": self.kb_id,
            "documents": [
                {
                    "document_id": document_id,
                    "status": "processed",
                    "chunks_count": 1,
                }
                for document_id in document_ids
            ],
        }

    def _query_data(self, query: str, *, answer: str | None = None) -> dict[str, Any]:
        references = [
            {"document_id": document_id, "file_path": item["source"]}
            for document_id, item in self.documents.items()
        ]
        chunks = [
            {
                "chunk_id": f"chunk_{document_id}",
                "document_id": document_id,
                "content": item["text"],
                "score": 1.0,
            }
            for document_id, item in self.documents.items()
        ]
        context = "\n".join(item["text"] for item in self.documents.values())
        return {
            "kb_id": self.kb_id,
            "hit": bool(chunks),
            "query": query,
            "effective_query": query,
            "rewritten": False,
            "has_context": bool(chunks),
            "context": context,
            "context_truncated": False,
            "answer": answer,
            "references": references,
            "chunks": chunks,
            "duration": 0.001,
        }

    async def query_context(
        self,
        query: str,
        profile: str,
        options: object,
        conversation: object,
    ) -> tuple[dict[str, Any], dict[str, Any]]:
        del profile, options, conversation
        return self._query_data(query), {"latency_ms": 1.0}

    async def query_answer(
        self,
        query: str,
        profile: str,
        options: object,
        conversation: object,
    ) -> tuple[dict[str, Any], dict[str, Any]]:
        del profile, options, conversation
        context = "\n".join(item["text"] for item in self.documents.values())
        return self._query_data(query, answer=context), {"latency_ms": 1.0}


class FakeHttpRagManager:
    """真实使用 SQLite 和 KB 目录，只替换在线 LightRAG Engine。"""

    def __init__(self, tmp_path: Path) -> None:
        self.metadata = MetadataStore(
            tmp_path / "liverag.db",
            tmp_path / "rag" / "knowledge_bases",
        )
        self.kb_store = KnowledgeBaseStore(self.metadata)
        self.kb_store.initialize()
        self.engines: dict[str, FakeHttpRagEngine] = {}

    async def get_engine(self, kb_id: str) -> FakeHttpRagEngine:
        self.kb_store.get(kb_id)
        return self.engines.setdefault(kb_id, FakeHttpRagEngine(kb_id))

    @staticmethod
    def _settings_for(meta: object) -> object:
        return {"kb_id": meta.kb_id, "kb_name": meta.name}


@pytest.fixture
def client() -> TestClient:
    """创建不会把服务端异常重新抛给测试进程的 HTTP 客户端。

    这里不使用上下文管理器，因此不会执行应用 lifespan，也就不会初始化
    真实的 LightRAG Engine；本文件只验证 HTTP 响应契约。
    """

    return TestClient(server.app, raise_server_exceptions=False)


@pytest.fixture
def isolated_http_client(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> TestClient:
    """创建使用临时 SQLite、隔离目录和 Fake Engine 的完整 HTTP 客户端。"""

    manager = FakeHttpRagManager(tmp_path)
    monkeypatch.setattr(server, "manager", manager)
    return TestClient(server.app, raise_server_exceptions=False)


def _auth_headers() -> dict[str, str]:
    """当测试环境启用了内部 API Key 时，返回合法认证头。"""

    if not server.settings.api_key:
        return {}
    return {"X-API-Key": server.settings.api_key}


def _assert_request_id(value: object) -> None:
    """断言 request_id 是非空且合法的 UUID 字符串。"""

    assert isinstance(value, str)
    assert str(UUID(value)) == value


def test_healthz_returns_success_envelope(client: TestClient) -> None:
    """成功响应包含统一 envelope 的全部字段。"""

    response = client.get("/v1/healthz")

    assert response.status_code == 200
    payload = response.json()

    assert set(payload) == {"request_id", "status", "data", "metrics", "error"}
    _assert_request_id(payload["request_id"])
    assert payload["status"] == "ok"
    assert payload["data"] == {"service": "ok"}
    assert payload["metrics"] == {}
    assert payload["error"] is None


def test_patch_knowledge_base_updates_metadata_and_returns_detail(
    isolated_http_client: TestClient,
) -> None:
    created = isolated_http_client.post(
        "/v1/knowledge-bases",
        headers=_auth_headers(),
        json={"name": "修改前", "description": "旧描述"},
    ).json()["data"]

    response = isolated_http_client.patch(
        f"/v1/knowledge-bases/{created['kb_id']}",
        headers=_auth_headers(),
        json={"name": "test", "description": "测试类"},
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["status"] == "ok"
    assert payload["data"]["kb_id"] == created["kb_id"]
    assert payload["data"]["name"] == "test"
    assert payload["data"]["description"] == "测试类"


def test_unhandled_exception_returns_error_envelope(
    client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """未处理异常被转换成不泄露内部细节的统一 500 envelope。"""

    def raise_unexpected_error() -> list[dict[str, object]]:
        raise RuntimeError("不应暴露给客户端的内部细节")

    monkeypatch.setattr(server.manager.kb_store, "list", raise_unexpected_error)

    response = client.get(
        "/v1/knowledge-bases",
        headers=_auth_headers(),
    )

    assert response.status_code == 500
    payload = response.json()

    assert set(payload) == {"request_id", "status", "data", "metrics", "error"}
    _assert_request_id(payload["request_id"])
    assert payload["status"] == "error"
    assert payload["data"] is None
    assert payload["metrics"] == {}
    assert payload["error"] == {
        "type": "InternalServerError",
        "message": "RAG Core 处理请求时发生内部错误",
    }
    assert "不应暴露给客户端的内部细节" not in response.text


def test_http_exception_returns_error_envelope(
    client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """已知 HTTP 错误也使用统一 envelope，并保留正确状态码。"""

    def raise_not_found(_: str) -> dict[str, object]:
        raise KeyError("知识库不存在")

    monkeypatch.setattr(server.manager.kb_store, "public_detail", raise_not_found)

    response = client.get(
        "/v1/knowledge-bases/missing",
        headers=_auth_headers(),
    )

    assert response.status_code == 404
    payload = response.json()
    assert set(payload) == {"request_id", "status", "data", "metrics", "error"}
    _assert_request_id(payload["request_id"])
    assert payload["status"] == "error"
    assert payload["data"] is None
    assert payload["metrics"] == {}
    assert payload["error"] == {
        "type": "HTTPException",
        "message": "知识库不存在",
    }


def test_request_validation_error_returns_error_envelope(client: TestClient) -> None:
    """缺少必填请求字段时返回统一的 422 envelope。"""

    response = client.post(
        "/v1/knowledge-bases/default/query/context",
        headers=_auth_headers(),
        json={},
    )

    assert response.status_code == 422
    payload = response.json()
    assert set(payload) == {"request_id", "status", "data", "metrics", "error"}
    _assert_request_id(payload["request_id"])
    assert payload["status"] == "error"
    assert payload["data"] is None
    assert payload["metrics"] == {}
    assert payload["error"]["type"] == "RequestValidationError"
    assert payload["error"]["message"] == "请求参数校验失败"
    assert payload["error"]["details"]


def test_write_source_file_keeps_target_inside_document_directory(
    tmp_path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """即使传入带目录的文件名，最终文件也只能写入文档目录。"""

    document_directory = tmp_path / "sources" / "doc_test"
    monkeypatch.setattr(
        server.manager.metadata,
        "source_document_dir",
        lambda **_: document_directory,
    )

    path = server._write_source_file(
        kb_id="default",
        document_id="doc_test",
        filename="../../outside.txt",
        raw=b"safe content",
    )

    assert path.parent == document_directory.resolve()
    assert path.name == "outside.txt"
    assert path.read_bytes() == b"safe content"
    assert not (tmp_path / "outside.txt").exists()


def test_http_rag_flow_keeps_two_knowledge_bases_isolated(
    isolated_http_client: TestClient,
) -> None:
    """通过 HTTP 完成建库、上传、job、context 和 answer，并验证两库隔离。"""

    client = isolated_http_client
    headers = _auth_headers()

    alpha_create = client.post(
        "/v1/knowledge-bases",
        headers=headers,
        json={"name": "Alpha", "description": "Alpha 测试库"},
    )
    beta_create = client.post(
        "/v1/knowledge-bases",
        headers=headers,
        json={"name": "Beta", "description": "Beta 测试库"},
    )
    assert alpha_create.status_code == 200
    assert beta_create.status_code == 200
    alpha_id = alpha_create.json()["data"]["kb_id"]
    beta_id = beta_create.json()["data"]["kb_id"]
    assert alpha_id != beta_id

    alpha_fact = "Alpha 项目的秘密代号是 ORANGE。"
    beta_fact = "Beta 项目的秘密代号是 PURPLE。"
    alpha_upload = client.post(
        f"/v1/knowledge-bases/{alpha_id}/documents/files",
        headers=headers,
        files={"files": ("alpha.txt", alpha_fact.encode(), "text/plain")},
    )
    beta_upload = client.post(
        f"/v1/knowledge-bases/{beta_id}/documents/files",
        headers=headers,
        files={"files": ("beta.md", beta_fact.encode(), "text/markdown")},
    )
    assert alpha_upload.status_code == 200
    assert beta_upload.status_code == 200
    alpha_upload_data = alpha_upload.json()["data"]
    beta_upload_data = beta_upload.json()["data"]
    assert alpha_upload_data["parsed_count"] == 1
    assert beta_upload_data["parsed_count"] == 1

    alpha_path = Path(alpha_upload_data["files"][0]["source_file_path"])
    beta_path = Path(beta_upload_data["files"][0]["source_file_path"])
    assert alpha_path.read_text(encoding="utf-8") == alpha_fact
    assert beta_path.read_text(encoding="utf-8") == beta_fact
    assert alpha_id in alpha_path.parts
    assert beta_id in beta_path.parts
    assert alpha_path.parent != beta_path.parent

    for kb_id, upload_data in (
        (alpha_id, alpha_upload_data),
        (beta_id, beta_upload_data),
    ):
        job_response = client.get(
            f"/v1/knowledge-bases/{kb_id}/jobs/{upload_data['track_id']}",
            headers=headers,
        )
        assert job_response.status_code == 200
        job_data = job_response.json()["data"]
        assert job_data["kb_id"] == kb_id
        assert job_data["status"] == "processed"
        assert job_data["documents"][0]["index_status"] == "processed"

    alpha_context = client.post(
        f"/v1/knowledge-bases/{alpha_id}/query/context",
        headers=headers,
        json={"query": "项目代号是什么？", "include_chunk_content": True},
    )
    beta_context = client.post(
        f"/v1/knowledge-bases/{beta_id}/query/context",
        headers=headers,
        json={"query": "项目代号是什么？", "include_chunk_content": True},
    )
    assert alpha_context.status_code == 200
    assert beta_context.status_code == 200
    alpha_result = alpha_context.json()["data"]
    beta_result = beta_context.json()["data"]
    assert isinstance(alpha_result, dict)
    assert isinstance(beta_result, dict)
    assert alpha_result["kb_id"] == alpha_id
    assert beta_result["kb_id"] == beta_id
    assert alpha_result["hit"] is True
    assert beta_result["hit"] is True
    assert alpha_result["duration"] >= 0
    assert beta_result["duration"] >= 0
    assert alpha_result["references"] and alpha_result["chunks"]
    assert beta_result["references"] and beta_result["chunks"]
    assert "ORANGE" in alpha_result["context"]
    assert "PURPLE" not in alpha_result["context"]
    assert "PURPLE" in beta_result["context"]
    assert "ORANGE" not in beta_result["context"]

    alpha_answer = client.post(
        f"/v1/knowledge-bases/{alpha_id}/query/answer",
        headers=headers,
        json={"query": "项目代号是什么？"},
    )
    beta_answer = client.post(
        f"/v1/knowledge-bases/{beta_id}/query/answer",
        headers=headers,
        json={"query": "项目代号是什么？"},
    )
    assert alpha_answer.status_code == 200
    assert beta_answer.status_code == 200
    alpha_answer_result = alpha_answer.json()["data"]
    beta_answer_result = beta_answer.json()["data"]
    assert isinstance(alpha_answer_result, dict)
    assert isinstance(beta_answer_result, dict)
    assert "ORANGE" in alpha_answer_result["answer"]
    assert "PURPLE" not in alpha_answer_result["answer"]
    assert "PURPLE" in beta_answer_result["answer"]
    assert "ORANGE" not in beta_answer_result["answer"]


def test_mixed_file_upload_preserves_failure_and_indexes_only_valid_document(
    isolated_http_client: TestClient,
) -> None:
    """同批文件部分失败时保留原件和错误，但只索引成功文档。"""

    response = isolated_http_client.post(
        "/v1/knowledge-bases/default/documents/files",
        headers=_auth_headers(),
        files=[
            ("files", ("valid.txt", b"valid searchable content", "text/plain")),
            (
                "files",
                (
                    "broken.docx",
                    b"not a valid office package",
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ),
            ),
        ],
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["status"] == "ok"
    data = payload["data"]
    assert data["parsed_count"] == 1
    assert data["error_count"] == 1
    assert data["total_files"] == 2

    documents = {item["original_filename"]: item for item in data["files"]}
    failed = documents["broken.docx"]
    successful = documents["valid.txt"]
    assert failed["status"] == "parse_failed"
    assert failed["error_msg"]
    assert Path(failed["source_file_path"]).read_bytes() == b"not a valid office package"
    assert Path(successful["source_file_path"]).read_bytes() == b"valid searchable content"

    manager = server.manager
    job = manager.metadata.job_detail("default", data["track_id"])
    assert job["parsed_count"] == 1
    assert job["failed_count"] == 1
    job_documents = {item["document_id"]: item for item in job["documents"]}
    assert job_documents[failed["document_id"]]["status"] == "parse_failed"
    assert job_documents[failed["document_id"]]["error_msg"]

    engine = manager.engines["default"]
    assert set(engine.documents) == {successful["document_id"]}
    assert failed["document_id"] not in engine.documents


@pytest.mark.parametrize(
    ("filename", "content_type", "raw", "expected_text"),
    [
        ("sample.pdf", "application/pdf", _pdf_bytes("PDF_FACT"), "PDF_FACT"),
        (
            "sample.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            _docx_bytes("DOCX_FACT"),
            "DOCX_FACT",
        ),
        (
            "sample.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            _pptx_bytes("PPTX_FACT"),
            "PPTX_FACT",
        ),
        (
            "sample.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            _xlsx_bytes("XLSX_FACT"),
            "XLSX_FACT",
        ),
    ],
    ids=["pdf", "docx", "pptx", "xlsx"],
)
def test_binary_format_upload_indexes_and_queries(
    isolated_http_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
    filename: str,
    content_type: str,
    raw: bytes,
    expected_text: str,
) -> None:
    """每种 M3-D 格式均可通过 HTTP 上传、索引、查看状态并查询。"""

    def unavailable_docling(_):
        raise ModuleNotFoundError("docling")

    monkeypatch.setattr(doc_parser, "_extract_pdf_with_docling", unavailable_docling)

    upload = isolated_http_client.post(
        "/v1/knowledge-bases/default/documents/files",
        headers=_auth_headers(),
        files={"files": (filename, raw, content_type)},
    )

    assert upload.status_code == 200
    upload_data = upload.json()["data"]
    assert upload_data["parsed_count"] == 1
    assert upload_data["error_count"] == 0
    document = upload_data["files"][0]
    assert document["source_file_exists"] is True

    job = isolated_http_client.get(
        f"/v1/knowledge-bases/default/jobs/{upload_data['track_id']}",
        headers=_auth_headers(),
    )
    assert job.status_code == 200
    assert job.json()["data"]["status"] == "processed"
    assert job.json()["data"]["documents"][0]["index_status"] == "processed"

    query = isolated_http_client.post(
        "/v1/knowledge-bases/default/query/context",
        headers=_auth_headers(),
        json={"query": "事实是什么？", "include_chunk_content": True},
    )
    assert query.status_code == 200
    assert expected_text in query.json()["data"]["context"]


def test_encrypted_pdf_upload_accepts_password(
    isolated_http_client: TestClient,
) -> None:
    """RAG Core multipart 字段会把正确密码交给 PDF 解析器。"""

    response = isolated_http_client.post(
        "/v1/knowledge-bases/default/documents/files",
        headers=_auth_headers(),
        files={
            "files": (
                "encrypted.pdf",
                _pdf_bytes("SECRET_PDF_FACT", password="correct-password"),
                "application/pdf",
            )
        },
        data={"pdf_password": "correct-password"},
    )

    assert response.status_code == 200
    data = response.json()["data"]
    assert data["parsed_count"] == 1
    assert data["error_count"] == 0
    engine = server.manager.engines["default"]
    assert "SECRET_PDF_FACT" in next(iter(engine.documents.values()))["text"]
